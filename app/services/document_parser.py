from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


class DocumentParseError(RuntimeError):
    """Raised when an uploaded file cannot be parsed into text."""


@dataclass
class ParsedDocument:
    text: str
    source_type: str
    unit_count: int


class DocumentParserService:
    def extract_text(self, file_path: Path) -> ParsedDocument:
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return self._parse_pdf(file_path)
        if suffix == ".docx":
            return self._parse_docx(file_path)
        if suffix == ".pptx":
            return self._parse_pptx(file_path)
        raise DocumentParseError("Unsupported file type. Please upload a PDF, DOCX, or PPTX file.")

    def _parse_pdf(self, file_path: Path) -> ParsedDocument:
        try:
            reader = PdfReader(str(file_path))
            texts = [page.extract_text() or "" for page in reader.pages]
        except Exception as exc:  # pragma: no cover - library exception surface
            raise DocumentParseError("The PDF could not be read. Try exporting it again with selectable text.") from exc

        text = self._normalize_text("\n".join(texts))
        if not text:
            raise DocumentParseError("The PDF did not contain readable text. If it is scanned, use OCR before uploading.")
        return ParsedDocument(text=text, source_type="PDF", unit_count=len(reader.pages))

    def _parse_docx(self, file_path: Path) -> ParsedDocument:
        try:
            document = Document(str(file_path))
        except Exception as exc:  # pragma: no cover - library exception surface
            raise DocumentParseError("The Word document could not be opened. Please try a different file.") from exc

        parts: list[str] = []
        parts.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)

        text = self._normalize_text("\n".join(parts))
        if not text:
            raise DocumentParseError("The Word document appears to be empty.")
        return ParsedDocument(text=text, source_type="DOCX", unit_count=max(len(document.paragraphs), 1))

    def _parse_pptx(self, file_path: Path) -> ParsedDocument:
        try:
            presentation = Presentation(str(file_path))
        except Exception as exc:  # pragma: no cover - library exception surface
            raise DocumentParseError("The PowerPoint file could not be opened. Please try a different file.") from exc

        slides_content: list[str] = []
        for slide in presentation.slides:
            slide_text: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_text.append(row_text)
            if slide_text:
                slides_content.append("\n".join(slide_text))

        text = self._normalize_text("\n\n".join(slides_content))
        if not text:
            raise DocumentParseError("The presentation did not contain readable slide text.")
        return ParsedDocument(text=text, source_type="PPTX", unit_count=len(presentation.slides))

    @staticmethod
    def _normalize_text(text: str) -> str:
        compact = text.replace("\x00", " ")
        compact = re.sub(r"\r\n?", "\n", compact)
        compact = re.sub(r"[ \t]+", " ", compact)
        compact = re.sub(r"\n{3,}", "\n\n", compact)
        return compact.strip()

